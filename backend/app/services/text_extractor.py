from __future__ import annotations

import logging
import tempfile
import os
from pathlib import Path

from docx import Document
from pptx import Presentation

logger = logging.getLogger(__name__)


# ── PDF strategies ──────────────────────────────────────────────

def _repair_split_numbers(text: str) -> str:
    """Merge digit sequences that were split by spaces during PDF conversion.

    PDF→DOCX often breaks long numbers like student IDs (e.g. 202300300254
    becomes "2023003002 54"). This merges them back.
    """
    import re
    # Merge: 10-digit prefix + space + 2-digit suffix → 12-digit ID
    text = re.sub(r'(\d{10})\s+(\d{2})\b', r'\1\2', text)
    # Merge: 4-9 digit prefix + space + 2-5 digit suffix
    text = re.sub(r'(\d{4,9})\s+(\d{2,5})\b', r'\1\2', text)
    return text


def _extract_pdf_via_docx(file_path: str) -> tuple[str, str | None]:
    """Convert PDF to DOCX then extract — best for structured forms with tables."""
    from pdf2docx import Converter

    tmp_path = None
    try:
        fd, tmp_path = tempfile.mkstemp(suffix=".docx")
        os.close(fd)

        cv = Converter(file_path)
        cv.convert(tmp_path)
        cv.close()

        text = _extract_docx_text(tmp_path)
        if not text.strip():
            return "", "pdf2docx produced empty document"

        # Repair artifacts from PDF→DOCX conversion
        text = _repair_split_numbers(text)

        return text, None
    except Exception as exc:
        return "", f"pdf2docx: {exc}"
    finally:
        if tmp_path:
            try:
                os.unlink(tmp_path)
            except OSError:
                pass


def _extract_pdf_text_pdfplumber(file_path: str) -> tuple[str, str | None]:
    """Extract PDF text using pdfplumber — fast, layout-aware."""
    import pdfplumber

    text_parts: list[str] = []
    with pdfplumber.open(file_path) as pdf:
        for page in pdf.pages:
            page_text = page.extract_text()
            if page_text and page_text.strip():
                text_parts.append(page_text.strip())

            tables = page.extract_tables()
            for table in tables:
                if not table:
                    continue
                rows: list[str] = []
                for row in table:
                    if row is None:
                        continue
                    cells = [str(c or "").strip() for c in row]
                    rows.append(" | ".join(cells))
                if rows:
                    text_parts.append("[表格]\n" + "\n".join(rows))

    if not text_parts:
        return "", "PDF appears to contain no extractable text (possibly scanned image)"
    return "\n\n".join(text_parts), None


def _extract_pdf_text_pypdf(file_path: str) -> tuple[str, str | None]:
    """Fallback PDF extraction using pypdf."""
    from pypdf import PdfReader

    reader = PdfReader(file_path)
    chunks: list[str] = []
    for page in reader.pages:
        text = page.extract_text()
        if text and text.strip():
            chunks.append(text.strip())
    text = "\n".join(chunks).strip()
    if not text:
        return "", "No text extracted from PDF"
    return text, None


def _extract_pdf_urls_via_pdfplumber(file_path: str) -> str:
    """Extract hyperlink URLs from PDF annotations using pdfplumber.

    pdf2docx often drops hyperlinks; this recovers them.
    """
    import re
    import pdfplumber

    urls: list[str] = []
    try:
        with pdfplumber.open(file_path) as pdf:
            for page in pdf.pages:
                # Extract visible text URLs
                page_text = page.extract_text() or ""
                found = re.findall(r"https?://[^\s]+", page_text)
                urls.extend(found)

                # Extract annotation/hyperlink URLs
                if hasattr(page, "hyperlinks") and page.hyperlinks:
                    for link in page.hyperlinks:
                        uri = link.get("uri") or ""
                        if uri and "http" in uri.lower():
                            urls.append(uri)
    except Exception:
        return ""

    seen: set[str] = set()
    unique: list[str] = []
    for url in urls:
        url = url.rstrip(")】〕,，。.;;")
        if url not in seen:
            seen.add(url)
            unique.append(url)
    return "\n".join(f"[超链接] {u}" for u in unique)


def extract_pdf_text(file_path: str) -> tuple[str, str | None]:
    """Extract text from PDF with multi-strategy fallback.

    Strategies:
    1. pdf2docx → DOCX extraction (preserves tables & label-value structure)
       + URL recovery from pdfplumber annotations
    2. pdfplumber (fast, layout-aware)
    3. pypdf (last resort)
    """
    errors: list[str] = []

    # Strategy 1: pdf2docx (best for structured forms) + pdfplumber URLs
    try:
        text, err = _extract_pdf_via_docx(file_path)
        if err is None and text.strip():
            # Recover URLs that pdf2docx lost (hyperlinks in PDF annotations)
            url_text = _extract_pdf_urls_via_pdfplumber(file_path)
            if url_text:
                text = text + "\n\n" + url_text
            return text, None
        if err:
            errors.append(err)
    except Exception as exc:
        errors.append(f"pdf2docx: {exc}")

    # Strategy 2: pdfplumber
    try:
        text, err = _extract_pdf_text_pdfplumber(file_path)
        if err is None and text.strip():
            return text, None
        if err:
            errors.append(err)
    except Exception as exc:
        errors.append(f"pdfplumber: {exc}")

    # Strategy 3: pypdf
    try:
        text, err = _extract_pdf_text_pypdf(file_path)
        if err is None and text.strip():
            return text, None
        if err:
            errors.append(err)
    except Exception as exc:
        errors.append(f"pypdf: {exc}")

    return "", "; ".join(errors) if errors else "All PDF extraction strategies failed"


# ── DOCX / PPTX ─────────────────────────────────────────────────

def _extract_docx_text(file_path: str) -> str:
    """Extract text from DOCX including paragraphs and tables."""
    doc = Document(file_path)
    chunks = [p.text for p in doc.paragraphs if p.text and p.text.strip()]
    for table in doc.tables:
        rows: list[str] = []
        for row in table.rows:
            values = [cell.text.strip() for cell in row.cells if cell.text.strip()]
            if values:
                rows.append(" | ".join(values))
        if rows:
            chunks.append("[表格]\n" + "\n".join(rows))
    return "\n".join(chunks).strip()


def _extract_pptx_text(file_path: str) -> str:
    """Extract text from PPTX including tables and grouped shapes."""
    presentation = Presentation(file_path)
    chunks: list[str] = []
    for slide_no, slide in enumerate(presentation.slides, start=1):
        slide_text: list[str] = []
        for shape in slide.shapes:
            if shape.shape_type == 6:  # MSO_SHAPE_TYPE.GROUP
                try:
                    for child in shape.shapes:
                        value = getattr(child, "text", "")
                        if value and value.strip():
                            slide_text.append(value.strip())
                except Exception:
                    pass
            if shape.has_table:
                try:
                    table = shape.table
                    for row in table.rows:
                        cells = [cell.text.strip() for cell in row.cells if cell.text.strip()]
                        if cells:
                            slide_text.append(" | ".join(cells))
                except Exception:
                    pass
            value = getattr(shape, "text", "")
            if value and value.strip():
                slide_text.append(value.strip())
        if slide_text:
            chunks.append(f"第{slide_no}页：" + "；".join(slide_text))
    return "\n".join(chunks).strip()


# ── Main entry point ────────────────────────────────────────────

def extract_text(file_path: str) -> tuple[str, str | None]:
    """Extract text content from a file. Returns (text, error).

    Supported formats: txt, md, log, csv, pdf, docx, pptx.
    PDF extraction tries pdf2docx, pdfplumber, then pypdf in order.
    """
    path = Path(file_path)
    suffix = path.suffix.lower().lstrip(".")

    try:
        # Plain text files
        if suffix in {"txt", "md", "log", "csv", "json", "xml", "html", "htm",
                       "py", "js", "ts", "vue", "css", "yaml", "yml", "toml"}:
            for encoding in ("utf-8", "gbk", "latin-1"):
                try:
                    text = path.read_text(encoding=encoding, errors="ignore")
                    return text.strip(), None
                except Exception:
                    continue
            return "", "unable to read text file"

        # PDF
        if suffix == "pdf":
            return extract_pdf_text(str(path))

        # DOCX
        if suffix == "docx":
            try:
                text = _extract_docx_text(str(path))
                return text, None
            except Exception as exc:
                return "", str(exc)

        # PPTX
        if suffix == "pptx":
            try:
                text = _extract_pptx_text(str(path))
                return text, None
            except Exception as exc:
                return "", str(exc)

        return "", f"unsupported file type: .{suffix}"
    except Exception as exc:
        logger.exception("Failed to extract text from %s", file_path)
        return "", str(exc)
